"""
parser.py — extraction pure.
Retourne (text_blocks, tables, images) depuis n'importe quel fichier supporté.

Formats supportés
─────────────────
.docx  → python-docx  (structural fidelity: multilingual headings, lists, tables,
                        images in document order; markitdown as fallback)
.pdf   → pypdf  (texte) + pdfplumber  (tableaux + colonnes) + Ollama/Gemma 4  (OCR scanned)
.pptx  → markitdown  (slide par slide, [SECTION] markers)
.xlsx  → openpyxl  (chaque feuille → schema chunk + data chunks)
.eml   → email parser  (exports Outlook; attachments dispatched by type)
.txt / .md / .csv / .json / .yaml / ...  → lecture texte brute

Double-column PDF support
─────────────────────────
pdfplumber est utilisé pour détecter et trier les blocs de texte par colonne.
Heuristique : si la page a des mots répartis sur deux zones X distinctes,
on trie colonne gauche d'abord, puis colonne droite (top→bottom dans chaque).
"""
from __future__ import annotations

import base64
import gc
import io
import os
import re
import sys
import tempfile
import zipfile
from contextlib import contextmanager
from email.parser import Parser as EmailParser
from pathlib import Path
from typing import Any, Callable, Generator, IO, Optional

import chardet
from loguru import logger
from PIL import Image

from file_processor.models import ExtractedImage, ExtractedTable


MAX_IMAGE_BYTES   = 20 * 1024 * 1024   # 20 MB
MIN_IMAGE_DIM     = 80                  # px — ignorer les micro-images

PLAIN_TEXT_EXTS   = {".txt", ".md", ".mdx",
                     ".json", ".xml", ".yml", ".yaml", ".sql", ".log", ".conf"}
CSV_EXTS          = {".csv", ".tsv"}
PASSWORD_EXTS     = {".pdf", ".docx", ".pptx", ".xlsx"}

# ── Double-column detection thresholds ───────────────────────────────────────
# Si le gap entre deux clusters X dépasse ce % de la largeur de page → double colonne
COLUMN_GAP_RATIO  = 0.10   # 10 % de la largeur de page
# Nombre minimal de mots pour tenter la détection de colonnes
MIN_WORDS_FOR_COL = 30

_MD = None
def _markitdown():
    global _MD
    if _MD is None:
        from markitdown import MarkItDown
        _MD = MarkItDown(enable_plugins=False)
    return _MD


# ── Helpers ───────────────────────────────────────────────────────────────────

@contextmanager
def _keep_pos(file: IO[Any]) -> Generator[IO[Any], None, None]:
    pos = file.tell()
    try:
        file.seek(0); yield file
    finally:
        file.seek(pos)


def _encoding(file: IO[bytes]) -> str:
    raw = file.read(50_000); file.seek(0)
    return chardet.detect(raw)["encoding"] or "utf-8"


def _b64(raw: bytes) -> str:
    return base64.b64encode(raw).decode()


def _resize(data: bytes) -> bytes:
    if len(data) <= MAX_IMAGE_BYTES:
        return data
    logger.info(f"  image redimensionnée ({len(data)//1_048_576} MB → 1024px)")
    with Image.open(io.BytesIO(data)) as img:
        img.thumbnail((1024, 1024), Image.Resampling.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return buf.getvalue()


def _make_image(raw: bytes, page: int, idx: int) -> Optional[ExtractedImage]:
    """Valide, redimensionne, encode → ExtractedImage  ou  None si trop petite."""
    try:
        pil = Image.open(io.BytesIO(raw))
        if pil.width < MIN_IMAGE_DIM or pil.height < MIN_IMAGE_DIM:
            return None
        return ExtractedImage(page_number=page, image_index=idx,
                              base64_data=_b64(_resize(raw)))
    except Exception as e:
        logger.warning(f"  image ignorée (page={page}, idx={idx}): {e}")
        return None


def _rows_to_md(rows: list[list[Any]]) -> str:
    if not rows:
        return ""
    sr = [[str(c or "").strip() for c in r] for r in rows]
    n  = max(len(r) for r in sr)
    sr = [r + [""] * (n - len(r)) for r in sr]
    h  = "| " + " | ".join(sr[0])  + " |"
    sep= "| " + " | ".join(["---"]*n) + " |"
    body = "\n".join("| " + " | ".join(r) + " |" for r in sr[1:])
    return "\n".join(filter(None, [h, sep, body]))


# ── Double-column detection & sorting ─────────────────────────────────────────

def _detect_column_split(words: list[dict], page_width: float) -> Optional[float]:
    """
    Détecte la frontière X entre deux colonnes.

    Retourne la coordonnée X du milieu du gap, ou None si page simple colonne.

    Stratégie :
    1. Construire un histogramme horizontal des centres de mots (x0+x1)/2
    2. Chercher un gap significatif dans la zone centrale (30%–70%) de la page
    3. Si ce gap est > COLUMN_GAP_RATIO * page_width → double colonne détectée
    """
    if not words or page_width <= 0:
        return None

    centers = sorted((w["x0"] + w["x1"]) / 2 for w in words)
    if len(centers) < MIN_WORDS_FOR_COL:
        return None

    # Zone centrale de la page (on ne cherche pas de colonne tout à gauche/droite)
    zone_left  = page_width * 0.30
    zone_right = page_width * 0.70
    center_pts = [c for c in centers if zone_left <= c <= zone_right]

    if len(center_pts) < 4:
        return None

    # Chercher le plus grand gap dans la zone centrale
    best_gap   = 0.0
    best_split = None
    for i in range(len(center_pts) - 1):
        gap = center_pts[i + 1] - center_pts[i]
        if gap > best_gap:
            best_gap   = gap
            best_split = (center_pts[i] + center_pts[i + 1]) / 2

    threshold = page_width * COLUMN_GAP_RATIO
    if best_gap >= threshold and best_split is not None:
        logger.debug(f"  double colonne détectée — split X={best_split:.1f} gap={best_gap:.1f}")
        return best_split

    return None


def _words_to_blocks_single(words: list[dict]) -> list[str]:
    """
    Regroupe les mots d'une page simple colonne en lignes de texte.
    Les mots sont triés par (top, x0) puis regroupés en lignes
    si leur chevauchement vertical est suffisant.
    """
    if not words:
        return []

    # Tolérance verticale pour appartenir à la même ligne
    tolerance = 3.0

    lines: list[list[dict]] = []
    for w in sorted(words, key=lambda w: (round(w["top"] / tolerance), w["x0"])):
        placed = False
        for line in lines:
            ref_top = line[0]["top"]
            if abs(w["top"] - ref_top) <= tolerance:
                line.append(w)
                placed = True
                break
        if not placed:
            lines.append([w])

    blocks = []
    for line in lines:
        text = " ".join(w["text"] for w in sorted(line, key=lambda w: w["x0"]))
        text = text.strip()
        if text:
            blocks.append(text)

    return blocks


def _words_to_blocks_double(words: list[dict], split_x: float) -> list[str]:
    """
    Regroupe les mots d'une page double colonne en lignes de texte.
    Trie : colonne gauche (top→bottom) puis colonne droite (top→bottom).
    """
    left_words  = [w for w in words if (w["x0"] + w["x1"]) / 2 < split_x]
    right_words = [w for w in words if (w["x0"] + w["x1"]) / 2 >= split_x]

    blocks = (
        _words_to_blocks_single(left_words) +
        _words_to_blocks_single(right_words)
    )
    return blocks


def _is_real_table(rows: list) -> bool:
    """
    Return True only if the extracted rows look like a genuine data table.

    Design-tool PDFs (Canva, Adobe Express, etc.) contain decorative grid lines
    that pdfplumber interprets as table borders, producing 'tables' where almost
    every cell is empty.  These are layout artefacts, not real tables.

    Heuristic: a real table has at least 20 % of its cells filled with content.
    Anything below that threshold is treated as a false positive and discarded,
    which also prevents the table bounding-box from masking the text that sits
    on top of the decorative grid.

    Additionally, a table must have at least 2 rows and 2 columns — single-row
    or single-column 'tables' are almost always mis-detected text boxes.
    """
    if not rows or len(rows) < 2:
        return False
    # Require at least 2 columns in the majority of rows
    col_counts = [len(r) for r in rows]
    if max(col_counts) < 2:
        return False
    total_cells = sum(len(r) for r in rows)
    if total_cells == 0:
        return False
    non_empty = sum(1 for r in rows for c in r if c and str(c).strip())
    fill_ratio = non_empty / total_cells
    return fill_ratio >= 0.20   # <20 % fill → treat as decorative grid artefact


def _extract_page_text_blocks(page) -> list[str]:
    """
    Extrait les blocs de texte d'une page pdfplumber en gérant
    les layouts simple et double colonne.

    Args:
        page : objet pdfplumber.Page

    Returns:
        Liste de chaînes (une par ligne de texte reconstituée)
    """
    try:
        words = page.extract_words(
            x_tolerance=1.5,   # tighter than default (3) — prevents LaTeX/tight-spacing PDFs
                               # from merging whole sentences into a single "word"
            y_tolerance=3,
            keep_blank_chars=False,
            use_text_flow=False,
        )
    except Exception as e:
        logger.warning(f"  extract_words échoué page {page.page_number}: {e}")
        # Fallback : extraction simple
        text = page.extract_text() or ""
        return [l.strip() for l in text.splitlines() if l.strip()]

    if not words:
        return []

    page_width = float(page.width) if page.width else 0.0
    split_x    = _detect_column_split(words, page_width)

    if split_x is not None:
        return _words_to_blocks_double(words, split_x)
    else:
        return _words_to_blocks_single(words)


# ── Lock detection ─────────────────────────────────────────────────────────────

def _pdf_locked(f: IO) -> bool:
    from pypdf import PdfReader
    with _keep_pos(f): return PdfReader(f).is_encrypted

def _office_locked(f: IO) -> bool:
    try:
        import msoffcrypto
        with _keep_pos(f): return msoffcrypto.OfficeFile(f).is_encrypted()
    except Exception: return False

_LOCK_CHECK: dict[str, Callable] = {
    ".pdf": _pdf_locked, ".docx": _office_locked,
    ".pptx": _office_locked, ".xlsx": _office_locked,
}


# ── Parsers ───────────────────────────────────────────────────────────────────

def _docx_heading_level(style_name: str) -> int | None:
    """
    Returns the heading level (1–6) for a Word paragraph style, or None if the
    paragraph is not a heading.

    Handles the most common style name conventions across languages:
      English : "Heading 1" … "Heading 6"
      French  : "Titre 1"   … "Titre 6"     (also "Titre" alone → 1)
      German  : "Überschrift 1" …
      Dutch   : "Kop 1" …
      Italian : "Intestazione 1" …
      Spanish : "Título 1" …  / "Encabezado 1" …
      Portuguese: "Título 1" … / "Cabeçalho 1" …
      Document title styles → level 1

    Falls back to None for all body / list / table styles.
    """
    if not style_name:
        return None
    sn = style_name.strip().lower()

    # Numbered heading prefixes (e.g. "heading 2", "titre 3")
    _HEADING_PREFIXES = (
        "heading ",       # EN
        "titre ",         # FR
        "überschrift ",   # DE
        "kop ",           # NL
        "intestazione ",  # IT
        "título ",        # ES / PT
        "encabezado ",    # ES
        "cabeçalho ",     # PT
        "rubrik ",        # SV
        "overskrift ",    # DA / NO
    )
    for prefix in _HEADING_PREFIXES:
        if sn.startswith(prefix):
            remainder = sn[len(prefix):].strip()
            # Handle "heading 1", "heading1", "titre 2 car" etc.
            m = re.match(r'^(\d+)', remainder)
            if m:
                lvl = int(m.group(1))
                return lvl if 1 <= lvl <= 6 else None

    # Unnumbered title styles → treat as H1
    if sn in ("title", "titre", "titel", "titolo", "título", "tittel"):
        return 1

    # Subtitle → H2
    if sn in ("subtitle", "sous-titre", "untertitel", "subtítulo"):
        return 2

    return None


# Bullet-like prefix patterns already embedded in the text by the renderer
_BULLET_START_RE = re.compile(r'^[•●○◦▪\-\*]|\d+[.)]\s')


def _is_docx_list_para(para) -> bool:
    """
    True when a python-docx Paragraph belongs to a numbered or bulleted list.

    Detection order (cheapest first):
      1. Style name contains "list" (covers "List Paragraph", "List Bullet", …)
      2. XML: w:pPr/w:numPr is present — Word's canonical list indicator
    """
    try:
        sn = (para.style.name or "").lower()
        if "list" in sn:
            return True
    except Exception:
        pass
    try:
        from docx.oxml.ns import qn as _qn
        pPr = para._p.find(_qn("w:pPr"))
        if pPr is not None and pPr.find(_qn("w:numPr")) is not None:
            return True
    except Exception:
        pass
    return False


_DOCX_CHARS_PER_PAGE: int = 3_000
"""
Heuristic: average characters per A4/Letter page of body text.
~500 words × 6 chars/word = 3 000 chars.  Overestimated slightly so
short documents don't balloon to many pages.  Resulting page_start values
are estimates — they may be off by ±1 page relative to a real renderer,
but are far more useful than page_start=1 for every chunk.
"""


def _parse_docx(path: Path):
    """
    Extraction DOCX avec python-docx :
    - Paragraphes → text blocks with heuristic page numbers
    - Tableaux → ExtractedTable objects
    - Images  → ExtractedImage depuis word/media/
    L'ordre document est préservé en itérant les enfants XML du body.

    Page estimation
    ---------------
    python-docx has no page-boundary awareness.  We accumulate character
    counts and increment the page estimate every _DOCX_CHARS_PER_PAGE
    characters (default 3 000 ≈ one A4 page of body text).  The estimate
    is imprecise — headings, list items, and tables consume fewer chars
    per visual page — but is far more useful for citations than page=1
    on every chunk.  The true page count is stored in
    doc_metadata["estimated_pages"].
    """
    data = path.read_bytes()
    blocks: list[tuple[int, str]] = []
    tables: list[ExtractedTable] = []
    images: list[ExtractedImage] = []
    img_idx = 0
    tbl_idx = 0

    # Page estimation state
    _char_accum: int = 0   # characters accumulated since page boundary
    _est_page:   int = 1   # current estimated page number (1-based)

    try:
        from docx import Document as DocxDocument
        from docx.table import Table as DocxTable
        from docx.text.paragraph import Paragraph as DocxPara

        # python-docx ne gère pas les chemins non-ASCII sur Windows → temp file
        import tempfile, os as _os
        tmp_fd2, tmp_path2 = tempfile.mkstemp(suffix=".docx")
        try:
            with _os.fdopen(tmp_fd2, "wb") as tmp_f2:
                tmp_f2.write(data)
            doc = DocxDocument(tmp_path2)
        finally:
            try: _os.unlink(tmp_path2)
            except Exception: pass

        # Buffer for consecutive list-item paragraphs.
        # Flushed into a single block so short items don't fragment into
        # dozens of tiny chunks that the chunker can't merge cleanly.
        _list_buf: list[str] = []

        def _flush_list_buf() -> None:
            nonlocal _char_accum, _est_page
            if _list_buf:
                merged = "\n".join(_list_buf)
                blocks.append((_est_page, merged))
                _char_accum += len(merged)
                if _char_accum >= _DOCX_CHARS_PER_PAGE:
                    _char_accum -= _DOCX_CHARS_PER_PAGE
                    _est_page += 1
                _list_buf.clear()

        # Track the nearest preceding heading for table context labelling.
        # DOCX has no page numbers so we can't use page-based context;
        # we must track heading state in document order at parse time.
        _current_heading: str | None = None

        for child in doc.element.body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                para = DocxPara(child, doc)
                text = para.text.strip()
                if not text:
                    continue
                try:
                    style_name = para.style.name or ""
                except Exception:
                    style_name = ""
                lvl = _docx_heading_level(style_name)
                if lvl:
                    # Headings always break any open list buffer
                    _flush_list_buf()
                    _current_heading = text   # plain text, no # prefix yet
                    text = f"{'#' * lvl} {text}"
                    blocks.append((_est_page, text))
                    _char_accum += len(text)
                    if _char_accum >= _DOCX_CHARS_PER_PAGE:
                        _char_accum -= _DOCX_CHARS_PER_PAGE
                        _est_page += 1
                    continue

                if _is_docx_list_para(para):
                    # Ensure a visible bullet prefix for numbered lists whose
                    # number is auto-generated (not embedded in the text run)
                    if not _BULLET_START_RE.match(text):
                        text = f"- {text}"
                    _list_buf.append(text)
                    _char_accum += len(text)
                else:
                    _flush_list_buf()
                    _char_accum += len(text)
                    blocks.append((_est_page, text))

                # Advance estimated page whenever we cross the threshold
                if _char_accum >= _DOCX_CHARS_PER_PAGE:
                    _char_accum -= _DOCX_CHARS_PER_PAGE
                    _est_page += 1
                continue   # already appended above (or into list buf)

            elif tag == "tbl":
                _flush_list_buf()
                tbl = DocxTable(child, doc)
                rows = []
                for row in tbl.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows.append(row_data)
                if rows:
                    tbl_chars = sum(
                        len(cell) for row in rows for cell in row
                    )
                    _char_accum += tbl_chars
                    if _char_accum >= _DOCX_CHARS_PER_PAGE:
                        _char_accum -= _DOCX_CHARS_PER_PAGE
                        _est_page += 1
                    tables.append(ExtractedTable(
                        page_number=_est_page,
                        table_index=tbl_idx,
                        markdown=_rows_to_md(rows),
                        raw_rows=rows,
                        section=_current_heading,
                    ))
                    tbl_idx += 1
                continue   # already handled

        # Flush any list items still buffered at end of document
        _flush_list_buf()

    except Exception as e:
        logger.warning(f"  python-docx échoué: {e} — fallback markitdown")
        try:
            from markitdown import StreamInfo
            res = _markitdown().convert(
                io.BytesIO(data),
                stream_info=StreamInfo(
                    mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                ),
            )
            md_text = res.markdown or ""
            blocks = [(0, b.strip()) for b in md_text.split("\n\n") if b.strip()]
        except Exception as e2:
            logger.warning(f"  markitdown aussi échoué: {e2}")

    # Images depuis le zip word/media/ — extraites en ordre document
    # Strategy:
    #   1. Parse word/_rels/document.xml.rels to build rId → media path map
    #   2. Walk word/document.xml body in order, collect r:embed rIds → ordered paths
    #   3. Fall back to sorted(namelist) if XML parsing fails
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            all_names = set(z.namelist())
            ordered_media: list[str] = []
            try:
                import xml.etree.ElementTree as _ET
                _NS_RELS  = "http://schemas.openxmlformats.org/package/2006/relationships"
                _NS_R     = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                _EMBED_ATTR = f"{{{_NS_R}}}embed"

                # rId → "word/media/imageN.xxx"
                rels_xml  = z.read("word/_rels/document.xml.rels")
                rels_root = _ET.fromstring(rels_xml)
                rid_to_target: dict[str, str] = {}
                for rel in rels_root.findall(f"{{{_NS_RELS}}}Relationship"):
                    rid    = rel.get("Id", "")
                    target = rel.get("Target", "")
                    if target.startswith("media/"):
                        rid_to_target[rid] = f"word/{target}"

                # Walk document body to find r:embed in document order
                doc_xml  = z.read("word/document.xml")
                doc_root = _ET.fromstring(doc_xml)
                seen: set[str] = set()
                for elem in doc_root.iter():
                    rid = elem.get(_EMBED_ATTR)
                    if rid and rid in rid_to_target:
                        path = rid_to_target[rid]
                        if path not in seen and path in all_names:
                            seen.add(path)
                            ordered_media.append(path)

                # Append any media files not referenced via r:embed
                # (e.g. background images, EMF/WMF objects) at the end
                for fallback in sorted(all_names):
                    if (fallback.startswith("word/media/")
                            and fallback not in seen):
                        ordered_media.append(fallback)

            except Exception as _xml_err:
                logger.debug(
                    f"  DOCX image order: XML parse failed ({_xml_err})"
                    " — falling back to sorted(namelist)"
                )
                ordered_media = [
                    n for n in sorted(all_names)
                    if n.startswith("word/media/")
                ]

            for name in ordered_media:
                img = _make_image(z.read(name), page=0, idx=img_idx)
                if img:
                    images.append(img)
                    img_idx += 1
    except Exception as e:
        logger.warning(f"  extraction images DOCX: {e}")

    # ── Métadonnées DOCX ─────────────────────────────────────────────────────
    doc_metadata: dict = {"file_type": "docx"}
    try:
        from docx import Document as _DocxMeta
        import tempfile, os as _os2
        tmp2_fd, tmp2_path = tempfile.mkstemp(suffix=".docx")
        try:
            with _os2.fdopen(tmp2_fd, "wb") as _f: _f.write(data)
            _dm = _DocxMeta(tmp2_path)
            cp = _dm.core_properties
            if cp.title:    doc_metadata["title"]            = cp.title
            if cp.author:   doc_metadata["author"]           = cp.author
            if cp.last_modified_by: doc_metadata["last_modified_by"] = cp.last_modified_by
            if cp.subject:  doc_metadata["subject"]          = cp.subject
            if cp.keywords: doc_metadata["keywords"]         = cp.keywords
            if cp.created:  doc_metadata["created"]          = str(cp.created)
            if cp.modified: doc_metadata["modified"]         = str(cp.modified)
            doc_metadata["paragraph_count"] = len(_dm.paragraphs)
        finally:
            try: _os2.unlink(tmp2_path)
            except Exception: pass
    except Exception as e:
        logger.debug(f"  DOCX metadata: {e}")

    # Store estimated page count so callers / chunker can surface it
    doc_metadata["estimated_pages"] = _est_page
    logger.info(
        f"  DOCX — {len(blocks)} blocs, {len(tables)} tableaux, {len(images)} images "
        f"(~{_est_page} page(s) estimated)"
    )
    return blocks, tables, images, doc_metadata


def _parse_pdf(path: Path, pdf_pass: Optional[str] = None):
    """
    Extraction PDF avec support double colonne.

    Étapes :
    1. pypdf  → images embarquées
    2. pdfplumber → texte avec détection de colonnes + tableaux
    3. OCR si PDF scanné (PaddleOCR)
    """
    from pypdf import PdfReader
    from pypdf.errors import PdfStreamError
    import pdfplumber

    blocks: list[str] = []
    tables: list[ExtractedTable] = []
    images: list[ExtractedImage] = []
    doc_metadata: dict = {"file_type": "pdf"}

    # ── Étape 1 : images via pypdf ────────────────────────────────────────────
    img_idx    = 0
    page_count = 0
    pypdf_text = ""

    try:
        reader = PdfReader(str(path))
        if reader.is_encrypted:
            if not pdf_pass or reader.decrypt(pdf_pass) == 0:
                logger.warning("  PDF chiffré — ignoré")
                return [], [], [], {"file_type": "pdf", "encrypted": True}

        page_count = len(reader.pages)
        doc_metadata["page_count"] = page_count

        # Extract PDF document metadata
        info = reader.metadata
        if info:
            if info.title:    doc_metadata["title"]     = info.title
            if info.author:   doc_metadata["author"]    = info.author
            if info.subject:  doc_metadata["subject"]   = info.subject
            if info.creator:  doc_metadata["creator"]   = info.creator
            if info.producer: doc_metadata["producer"]  = info.producer
            if info.creation_date:
                doc_metadata["created"]  = str(info.creation_date)
            if info.modification_date:
                doc_metadata["modified"] = str(info.modification_date)

        for pnum, page in enumerate(reader.pages):
            pypdf_text += page.extract_text() or ""
            for obj in page.images:
                img = _make_image(_resize(obj.data), page=pnum, idx=img_idx)
                if img:
                    images.append(img)
                    img_idx += 1

    except (PdfStreamError, Exception) as e:
        logger.error(f"  pypdf: {e}")

    # ── Étape 2 : détection PDF scanné → classification page par page ────────
    from file_processor.ocr    import is_scanned_pdf, POPPLER_PATH
    from file_processor.gemma4 import (
        gemma4_available,
        openrouter_available as _openrouter_available,
        ocr_image_bytes_gemma as _ocr_image_bytes,   # OpenRouter-first, Ollama fallback
        OCRMode as _GemmaOCRMode,
    )

    if is_scanned_pdf(pypdf_text, page_count):
        logger.info("  PDF scanné détecté — classification page par page ...")
        doc_metadata["scanned"] = True   # consumed by chunker to select TXT pipeline
        images = []

        _ocr_ok = _openrouter_available() or gemma4_available()
        if not _ocr_ok:
            logger.warning(
                "  OCR non disponible — configure OPENROUTER_API_KEY "
                "ou lance: ollama pull gemma4:e4b"
            )
        else:
            try:
                from pdf2image import convert_from_path
                pil_pages = convert_from_path(str(path), dpi=200, poppler_path=POPPLER_PATH)
                img_idx = 0

                _ocr_disabled = False  # set True after an OOM so we stop trying
                for pnum, pil_page in enumerate(pil_pages):
                    logger.info(f"  Analyse page {pnum + 1}/{len(pil_pages)} ...")

                    # Convert PIL page to bytes once — reused for both OCR and image fallback.
                    page_buf = io.BytesIO()
                    pil_page.save(page_buf, format="PNG")
                    page_bytes = page_buf.getvalue()

                    page_text = ""
                    if not _ocr_disabled:
                        try:
                            page_text = _ocr_image_bytes(page_bytes, mode=_GemmaOCRMode.EXTRACT)
                        except Exception as page_err:
                            err_str = str(page_err).lower()
                            if "memory" in err_str or "oom" in err_str or "out of memory" in err_str:
                                logger.warning(
                                    f"  → page {pnum + 1}: OCR OOM — "
                                    "OCR disabled for remaining pages. "
                                    "Set OPENROUTER_API_KEY or use a smaller Ollama model."
                                )
                                _ocr_disabled = True
                            else:
                                logger.warning(f"  → page {pnum + 1}: OCR error — {page_err}")

                    text_chars = len(page_text.strip())

                    if text_chars >= 50:
                        # ── Page texte → whole-page block (preserves OCR markdown) ──
                        # The OCR backend returns structured markdown (headings, tables,
                        # paragraphs).  Storing the whole page as one block lets the
                        # chunker reassemble table rows and heading hierarchy correctly —
                        # splitting by line would break markdown table structure.
                        _ocr_backend = "OpenRouter" if _openrouter_available() else "Gemma4/Ollama"
                        blocks.append((pnum, page_text.strip()))
                        logger.info(
                            f"  → page {pnum + 1}: {_ocr_backend} OCR ({text_chars} chars)"
                        )
                    else:
                        # ── Page visuelle → image chunk (no readable text) ────
                        img = _make_image(page_bytes, page=pnum, idx=img_idx)
                        if img:
                            images.append(img)
                            img_idx += 1
                        logger.info(
                            f"  → page {pnum + 1}: "
                            + ("skipped (OCR disabled)" if _ocr_disabled else f"image (peu de texte: {text_chars} chars)")
                        )

                logger.info(
                    f"  Terminé — {len(blocks)} blocs OCR (pages texte), "
                    f"{len(images)} images LLaVA (pages visuelles)"
                )

            except ImportError:
                logger.warning("  pdf2image non installé — lance: pip install pdf2image")
            except Exception as e:
                logger.error(f"  OCR/classification échoué : {e}")

        doc_metadata["is_scanned"] = True
        _extract_plumber_tables_only(path, tables)
        gc.collect()
        return blocks, tables, images, doc_metadata

    # ── Étape 3 : texte + colonnes via pdfplumber ─────────────────────────────
    logger.info(f"  PDF natif — extraction avec détection colonnes ...")
    try:
        with pdfplumber.open(str(path)) as pdf:
            for pnum, page in enumerate(pdf.pages):

                # 3a. Zones de tableaux — on les exclut du texte
                # Only real tables (≥20 % cell fill, ≥2×2) are kept.
                # Design-tool PDFs (Canva, etc.) have decorative grid lines that
                # pdfplumber misdetects as table borders — _is_real_table() filters
                # those out so their bboxes don't accidentally mask the real text.
                table_bboxes = []
                try:
                    page_tables = page.find_tables()
                    for tidx, tbl_obj in enumerate(page_tables):
                        rows = tbl_obj.extract()
                        if not rows:
                            continue
                        if not _is_real_table(rows):
                            logger.debug(
                                f"  p{pnum} table {tidx}: skipped (sparse/decorative "
                                f"{sum(1 for r in rows for c in r if c and str(c).strip())}/"
                                f"{sum(len(r) for r in rows)} cells filled)"
                            )
                            continue
                        table_bboxes.append(tbl_obj.bbox)
                        tables.append(ExtractedTable(
                            page_number=pnum,
                            table_index=tidx,
                            markdown=_rows_to_md(rows),
                            raw_rows=[[str(c or "") for c in r] for r in rows],
                        ))
                except Exception as e:
                    logger.warning(f"  pdfplumber tableaux p{pnum}: {e}")

                # 3b. Texte hors tableaux
                try:
                    if table_bboxes:
                        # Masquer les zones de tableaux avant extraction texte
                        text_page = page
                        for bbox in table_bboxes:
                            try:
                                text_page = text_page.filter(
                                    lambda obj, b=bbox: not (
                                        obj.get("x0", 0) >= b[0] and
                                        obj.get("x1", 0) <= b[2] and
                                        obj.get("top", 0) >= b[1] and
                                        obj.get("bottom", 0) <= b[3]
                                    )
                                )
                            except Exception:
                                pass
                        page_blocks = _extract_page_text_blocks(text_page)
                    else:
                        page_blocks = _extract_page_text_blocks(page)

                    blocks.extend([(pnum, b) for b in page_blocks])

                except Exception as e:
                    logger.warning(f"  extraction texte p{pnum}: {e}")
                    # Fallback pypdf pour cette page
                    page_text = (reader.pages[pnum].extract_text() or "") if page_count > pnum else ""
                    for line in page_text.splitlines():
                        if line.strip():
                            blocks.append((pnum, line.strip()))

    except Exception as e:
        logger.error(f"  pdfplumber: {e}")
        # Fallback complet sur pypdf
        for line in pypdf_text.splitlines():
            if line.strip():
                blocks.append((0, line.strip()))

    doc_metadata["is_scanned"] = False
    logger.info(f"  Extraction terminée — {len(blocks)} blocs, {len(tables)} tableaux, {len(images)} images")
    gc.collect()
    return blocks, tables, images, doc_metadata


def _extract_plumber_tables_only(path: Path, tables: list[ExtractedTable]) -> None:
    """Extrait uniquement les tableaux d'un PDF (pour les PDFs scannés)."""
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            for pnum, page in enumerate(pdf.pages):
                for tidx, tbl in enumerate(page.extract_tables() or []):
                    if tbl and _is_real_table(tbl):
                        tables.append(ExtractedTable(
                            page_number=pnum, table_index=tidx,
                            markdown=_rows_to_md(tbl),
                            raw_rows=[[str(c or "") for c in r] for r in tbl],
                        ))
    except Exception as e:
        logger.warning(f"  pdfplumber tableaux (scan): {e}")


def _parse_pptx(path: Path):
    """
    Extraction PPTX avec python-pptx :
    - Texte      → text blocks  (page = index de slide)
    - Tableaux   → ExtractedTable objects (page = index de slide)
    - Images     → ExtractedImage objects (page = index de slide)
    Chaque slide = une "page".
    """
    data = path.read_bytes()
    blocks: list[tuple[int, str]] = []
    tables: list[ExtractedTable] = []
    images: list[ExtractedImage] = []
    img_idx = 0
    tbl_idx = 0

    def _process_pptx_from_path(pptx_src):
        """Parse toutes les slides depuis un chemin ou un BytesIO."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        nonlocal img_idx, tbl_idx

        prs = Presentation(pptx_src)
        for slide_idx, slide in enumerate(prs.slides):

            # ── Detect section-divider slides ────────────────────────────────
            # Section divider slides use a layout whose name contains "section"
            # (e.g. "Section Header") AND have no TEXT_BOX shapes with real
            # content.  A slide that uses SECTION_HEADER layout but also has a
            # content TEXT_BOX (e.g. an intro/context slide) is NOT a divider
            # and must be treated as a normal content slide.
            try:
                layout_name = slide.slide_layout.name or ""
            except Exception:
                layout_name = ""

            has_content_textbox = any(
                sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
                and sh.has_text_frame
                and len(sh.text_frame.text.strip()) > 20
                for sh in slide.shapes
            )

            if "section" in layout_name.lower() and not has_content_textbox:
                # True section divider: collect only placeholder text (no
                # TEXT_BOX) to build the section label, skipping:
                #   • ph_idx >= 10  (slide number / date / footer)
                #   • pure-digit strings  (section numbers like "01")
                #   • the literal "‹#›" slide-number field
                section_texts: list[str] = []
                for sh in slide.shapes:
                    if not sh.has_text_frame:
                        continue
                    try:
                        ph_idx_sh = sh.placeholder_format.idx if sh.is_placeholder else None
                    except Exception:
                        ph_idx_sh = None
                    if ph_idx_sh is not None and ph_idx_sh >= 10:
                        continue  # skip slide-number / date / footer placeholders
                    for para in sh.text_frame.paragraphs:
                        t = para.text.strip()
                        if t and not re.fullmatch(r'\d+', t) and t != '‹#›':
                            section_texts.append(t)

                if section_texts:
                    label = " — ".join(section_texts)
                    blocks.append((slide_idx, f"[SECTION] {label}"))
                continue   # skip normal shape processing for section slides

            # ── Normal content slide ─────────────────────────────────────────
            # Strategy:
            #   • Emit title placeholder (ph_idx=0) as a markdown "# heading"
            #     so the chunker can use it as heading prefix.
            #   • Collect ALL other text (from body placeholders and TEXT_BOXes)
            #     into one combined body block per slide.  This prevents the
            #     common PPTX pattern of many tiny per-shape blocks that each
            #     get mis-classified as headings.
            #   • Skip layout/footer placeholders: ph_idx >= 10 (slide number,
            #     date, footer) and text consisting only of "‹#›".
            #   • Tables and images are emitted separately (not joined).
            #   • Title-only fallback: if nothing but the title placeholder was
            #     found, emit the title as a plain body sentence so the slide
            #     is never completely empty.

            slide_title_text:   str | None = None
            slide_has_table:    bool       = False
            slide_body_lines:   list[str]  = []   # accumulate all body text here

            for shape in slide.shapes:

                # ── Text ────────────────────────────────────────────────────
                if shape.has_text_frame:
                    # Identify placeholder index (None for TEXT_BOX)
                    try:
                        ph_idx = (
                            shape.placeholder_format.idx
                            if shape.is_placeholder else None
                        )
                    except Exception:
                        ph_idx = None

                    # Skip slide-number / date / footer placeholders
                    if ph_idx is not None and ph_idx >= 10:
                        continue

                    lines = [
                        para.text.strip()
                        for para in shape.text_frame.paragraphs
                        if para.text.strip() and para.text.strip() != '‹#›'
                    ]
                    if not lines:
                        continue

                    # Title placeholder (ph_idx == 0) → markdown heading block
                    if ph_idx == 0:
                        slide_title_text = lines[0]
                        blocks.append((slide_idx, f"# {lines[0]}"))
                        # Extra lines inside the title shape go to body
                        if len(lines) > 1:
                            slide_body_lines.extend(lines[1:])
                    else:
                        # All other text (body placeholders + TEXT_BOXes)
                        # accumulates into the per-slide body pool
                        slide_body_lines.extend(lines)

                # ── Table ────────────────────────────────────────────────────
                elif shape.has_table:
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    if rows:
                        tables.append(ExtractedTable(
                            page_number=slide_idx,
                            table_index=tbl_idx,
                            markdown=_rows_to_md(rows),
                            raw_rows=rows,
                        ))
                        tbl_idx += 1
                        slide_has_table = True

                # ── Image ────────────────────────────────────────────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = _make_image(shape.image.blob, page=slide_idx, idx=img_idx)
                        if img:
                            images.append(img)
                            img_idx += 1
                    except Exception as e:
                        logger.warning(f"  image slide {slide_idx}: {e}")

            # ── Emit combined body block ──────────────────────────────────────
            # One block per slide keeps related content together and avoids
            # many tiny per-shape blocks being mis-classified as headings.
            if slide_body_lines:
                blocks.append((slide_idx, "\n".join(slide_body_lines)))

            # ── Title-only fallback ───────────────────────────────────────────
            # If the slide has only a title (no body text, table, or image),
            # emit the title as a plain body sentence so the chunker never
            # skips the slide entirely.
            # "Slide: Title." — terminal period prevents heading re-detection.
            elif (slide_title_text
                      and not slide_has_table
                      and not any(img.page_number == slide_idx for img in images)):
                blocks.append((slide_idx, f"Slide: {slide_title_text}."))

    try:
        _process_pptx_from_path(io.BytesIO(data))

    except Exception as e:
        logger.warning(f"  python-pptx échoué: {e} — fallback markitdown + zip images")

        # ── Texte via markitdown ─────────────────────────────────────────────
        try:
            from markitdown import StreamInfo
            res = _markitdown().convert(
                io.BytesIO(data),
                stream_info=StreamInfo(
                    mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    filename=path.name, extension=".pptx",
                ),
            )
            md_text = res.markdown or ""
            blocks = [(0, b.strip()) for b in md_text.split("\n\n") if b.strip()]
        except Exception as e2:
            logger.warning(f"  markitdown aussi échoué: {e2}")

        # ── Images via zip direct (même si python-pptx échoue) ───────────────
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                for name in sorted(z.namelist()):
                    if name.startswith("ppt/media/"):
                        img = _make_image(z.read(name), page=0, idx=img_idx)
                        if img:
                            images.append(img)
                            img_idx += 1
            if images:
                logger.info(f"  {len(images)} image(s) extraite(s) via zip")
        except Exception as e3:
            logger.warning(f"  extraction images zip PPTX: {e3}")

    # ── Métadonnées PPTX ─────────────────────────────────────────────────────
    doc_metadata: dict = {"file_type": "pptx"}
    try:
        from pptx import Presentation as _PrsMeta
        _prs = _PrsMeta(io.BytesIO(data))
        cp = _prs.core_properties
        if cp.title:    doc_metadata["title"]    = cp.title
        if cp.author:   doc_metadata["author"]   = cp.author
        if cp.subject:  doc_metadata["subject"]  = cp.subject
        if cp.keywords: doc_metadata["keywords"] = cp.keywords
        if cp.created:  doc_metadata["created"]  = str(cp.created)
        if cp.modified: doc_metadata["modified"] = str(cp.modified)
        if cp.last_modified_by: doc_metadata["last_modified_by"] = cp.last_modified_by
        doc_metadata["slide_count"] = len(_prs.slides)
    except Exception as e:
        logger.debug(f"  PPTX metadata: {e}")

    logger.info(f"  PPTX — {len(blocks)} blocs, {len(tables)} tableaux, {len(images)} images")
    return blocks, tables, images, doc_metadata


def _parse_xlsx(path: Path):
    import openpyxl
    from zipfile import BadZipFile

    BUGS = ["Value must be either numerical",
            "File contains no valid workbook part",
            "Unable to read workbook"]
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except BadZipFile as e:
        logger.warning(f"  XLSX bad zip: {e}"); return [], [], []
    except Exception as e:
        if any(s in str(e) for s in BUGS):
            logger.error(f"  XLSX bug openpyxl: {e}"); return [], [], []
        raise

    BATCH_SIZE = 50   # lignes max par chunk (header répété dans chaque batch)

    blocks, tables = [], []
    tbl_global_idx = 0

    for sidx, ws in enumerate(wb.worksheets):
        rows, empty = [], 0
        for row in ws.iter_rows(min_row=1, values_only=True):
            parts  = [str(c or "").strip() for c in row]
            joined = ",".join(parts)
            if len(joined) >= len(row):
                rows.append(parts); empty = 0
            else:
                empty += 1
            if empty > 100:
                logger.warning(f"  feuille '{ws.title}': 100+ lignes vides — arrêt")
                break

        if not rows:
            continue

        header    = rows[0]
        data_rows = rows[1:]
        n_batches = max(1, (len(data_rows) + BATCH_SIZE - 1) // BATCH_SIZE)

        # Résumé texte de la feuille
        cols = ", ".join(header[:8]) + ("…" if len(header) > 8 else "")
        blocks.append((sidx,
            f"Feuille '{ws.title}': {len(data_rows)} lignes, "
            f"{len(header)} colonnes. Colonnes : {cols}."
        ))

        if n_batches == 1:
            # Petite feuille → un seul chunk table
            sheet_prefix = f"Sheet: {ws.title}\n"
            tables.append(ExtractedTable(
                page_number=sidx, table_index=tbl_global_idx,
                markdown=sheet_prefix + _rows_to_md(rows), raw_rows=rows,
            ))
            tbl_global_idx += 1
        else:
            # Grande feuille → batches de BATCH_SIZE lignes (header répété)
            logger.info(f"  feuille '{ws.title}': {len(data_rows)} lignes "
                        f"→ {n_batches} batches de {BATCH_SIZE}")
            for b in range(n_batches):
                batch_data = data_rows[b * BATCH_SIZE:(b + 1) * BATCH_SIZE]
                batch_rows = [header] + batch_data
                start_row  = b * BATCH_SIZE + 2        # 1-indexed, skip header
                end_row    = start_row + len(batch_data) - 1
                sheet_prefix = f"Sheet: {ws.title}, rows {start_row}–{end_row}\n"
                tables.append(ExtractedTable(
                    page_number=sidx, table_index=tbl_global_idx,
                    markdown=sheet_prefix + _rows_to_md(batch_rows), raw_rows=batch_rows,
                ))
                tbl_global_idx += 1

    # ── Métadonnées XLSX ─────────────────────────────────────────────────────
    doc_metadata: dict = {"file_type": "xlsx"}
    try:
        props = wb.properties
        if props.title:    doc_metadata["title"]    = props.title
        if props.creator:  doc_metadata["author"]   = props.creator
        if props.subject:  doc_metadata["subject"]  = props.subject
        if props.keywords: doc_metadata["keywords"] = props.keywords
        if props.created:  doc_metadata["created"]  = str(props.created)
        if props.modified: doc_metadata["modified"] = str(props.modified)
        if props.lastModifiedBy: doc_metadata["last_modified_by"] = props.lastModifiedBy
        doc_metadata["sheet_count"] = len(wb.worksheets)
        doc_metadata["sheet_names"] = [ws.title for ws in wb.worksheets]
    except Exception as e:
        logger.debug(f"  XLSX metadata: {e}")

    wb.close()   # release file handle so Windows can delete the file
    return blocks, tables, [], doc_metadata


def _parse_eml(path: Path):
    """
    Extraction EML :
    - En-têtes (Subject, From, To, Date, Attachments) → doc_metadata (pas dans le texte)
    - Corps text/plain  → text blocks
    - Corps text/html   → fallback si pas de text/plain (balises HTML supprimées)
    - Images inline (CID) → ExtractedImage
    Retourne (blocks, tables, images, doc_metadata).
    """
    import email as _email
    import email.policy as _policy

    def _strip_html(html: str) -> str:
        import re
        text = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.S | re.I)
        text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.S | re.I)
        text = re.sub(r'<[^>]+>', ' ', text)
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    def _extract_html_tables(html: str, out_tables: list, tbl_offset: int = 0) -> str:
        """
        Extract <table> elements from HTML into ExtractedTable objects.

        Each table is appended to out_tables and its HTML is replaced with a
        blank line so the surrounding prose remains intact.  Returns the
        modified HTML string (without the table elements).

        Falls back silently if BeautifulSoup is not available — tables stay
        embedded in the HTML and will be stripped with the rest of the tags.
        """
        try:
            from bs4 import BeautifulSoup, Tag
        except ImportError:
            logger.debug("  BeautifulSoup not available — HTML tables will not be extracted.")
            return html

        soup = BeautifulSoup(html, "html.parser")

        for tbl_idx, table_el in enumerate(soup.find_all("table")):
            rows: list[list[str]] = []
            for tr in table_el.find_all("tr"):
                cells = [
                    td.get_text(" ", strip=True)
                    for td in tr.find_all(["td", "th"])
                ]
                if any(c.strip() for c in cells):
                    rows.append(cells)

            if not rows:
                table_el.decompose()
                continue

            # Build markdown representation
            max_cols  = max(len(r) for r in rows)
            norm_rows = [r + [""] * (max_cols - len(r)) for r in rows]   # pad short rows

            md_header = "| " + " | ".join(norm_rows[0]) + " |"
            md_sep    = "| " + " | ".join(["---"] * max_cols) + " |"
            md_data   = "\n".join("| " + " | ".join(r) + " |" for r in norm_rows[1:])
            markdown  = "\n".join(filter(None, [md_header, md_sep, md_data]))

            out_tables.append(ExtractedTable(
                page_number = 0,
                table_index = tbl_offset + tbl_idx,
                markdown    = markdown,
                raw_rows    = norm_rows,
                section     = "email_body",
            ))

            # Replace table element with whitespace so surrounding prose is clean
            table_el.replace_with("\n\n")

        return str(soup)

    with open(path, "rb") as f:
        msg = _email.message_from_binary_file(f, policy=_policy.compat32)

    blocks: list[tuple[int, str]] = []
    tables: list[ExtractedTable] = []
    images: list[ExtractedImage] = []
    img_idx = 0

    # ── Métadonnées email → doc_metadata (séparées du contenu) ───────────────
    doc_metadata: dict = {}
    if msg.get("Subject"): doc_metadata["subject"]  = msg.get("Subject", "").strip()
    if msg.get("From"):    doc_metadata["from"]      = msg.get("From",    "").strip()
    if msg.get("To"):      doc_metadata["to"]        = msg.get("To",      "").strip()
    if msg.get("Cc"):      doc_metadata["cc"]        = msg.get("Cc",      "").strip()
    if msg.get("Date"):    doc_metadata["date"]      = msg.get("Date",    "").strip()
    if msg.get("Message-ID"): doc_metadata["message_id"] = msg.get("Message-ID", "").strip()

    # ── Corps et pièces jointes ───────────────────────────────────────────────
    plain_blocks:   list[str] = []
    html_blocks:    list[str] = []
    attachments:    list[str] = []
    html_had_tables: bool     = False   # True when HTML body contained ≥1 <table>

    def _is_likely_garbage(line: str) -> bool:
        """
        Return True for high-entropy random-character strings that typically
        appear as garbled/corrupted content in Outlook EML bodies.

        Heuristic: a token with no whitespace, length > 15, and a vowel ratio
        below 18 % is almost certainly not natural-language text.
        (Normal English: ~40 % vowels; typical garbage: < 10 %)
        """
        s = line.strip()
        if len(s) < 16 or ' ' in s or '\t' in s:
            return False
        alpha = [c for c in s if c.isalpha()]
        if len(alpha) < 10:
            return False
        vowel_ratio = sum(1 for c in alpha if c.lower() in "aeiou") / len(alpha)
        return vowel_ratio < 0.18

    for part in msg.walk():
        ctype       = part.get_content_type()
        disposition = str(part.get("Content-Disposition", ""))
        cid         = part.get("Content-ID", "")
        is_inline   = "inline" in disposition or cid
        is_attached = "attachment" in disposition

        # Text/plain
        if ctype == "text/plain" and not is_attached:
            payload = part.get_payload(decode=True)
            if payload:
                enc = part.get_content_charset() or "utf-8"
                text = payload.decode(enc, errors="replace")
                # ── Normalise line endings ─────────────────────────────────────
                # MIME bodies use CRLF (\r\n).  Normalise to \n so that
                # \n\n paragraph splitting and signature detection work correctly.
                text = text.replace("\r\n", "\n").replace("\r", "\n")
                # ── Strip email signature (RFC 3676 §4.3) ─────────────────────
                # A line containing exactly "-- " or "--" marks the start of the
                # signature block.  Everything from that line onward is discarded
                # before chunking — it adds no retrieval value and pollutes the
                # embedding space with names, phone numbers and tracking footers.
                import re as _re
                text = _re.split(r'(?m)^--\s*$', text, maxsplit=1)[0]
                # ── Strip inline CID image placeholders ───────────────────────
                # Outlook inserts [cid:...] markers into the plain-text body
                # for every inline image.  They carry no textual meaning and
                # would end up as literal tokens in the embedded chunks.
                text = _re.sub(r'\[cid:[^\]]+\]', '', text)
                # ── Filter garbage / high-entropy lines ───────────────────────
                # Some Outlook EML files include corrupt or randomly-encoded
                # lines in the plain-text part (artefacts of encoding issues).
                # Drop them line by line so surrounding real content survives.
                clean_lines = [
                    ln for ln in text.splitlines()
                    if not _is_likely_garbage(ln)
                ]
                text = "\n".join(clean_lines).strip()
                if text:
                    plain_blocks.append(text)

        # Text/html — always parse for tables; body text used when no plain or
        # when HTML contained tables (HTML body is cleaner in that case because
        # <table> elements are replaced with whitespace, avoiding the duplicate
        # one-cell-per-line table dump that Outlook puts in text/plain).
        elif ctype == "text/html" and not is_attached:
            payload = part.get_payload(decode=True)
            if payload:
                enc = part.get_content_charset() or "utf-8"
                html = payload.decode(enc, errors="replace")
                # Extract <table> elements as ExtractedTable before stripping HTML.
                # Tables are removed from the HTML so they don't appear as garbled
                # space-separated text in the body blocks.
                n_tables_before = len(tables)
                html = _extract_html_tables(html, tables, tbl_offset=len(tables))
                if len(tables) > n_tables_before:
                    html_had_tables = True
                cleaned = _strip_html(html)
                if cleaned:
                    html_blocks.append(cleaned)

        # Images inline (CID ou inline disposition)
        elif ctype.startswith("image/") and is_inline:
            payload = part.get_payload(decode=True)
            if payload:
                img = _make_image(payload, page=0, idx=img_idx)
                if img:
                    images.append(img)
                    img_idx += 1

        # Pièces jointes → extraire le contenu si le format est supporté
        elif is_attached:
            fname = part.get_filename() or ctype
            payload = part.get_payload(decode=True)
            ext = Path(fname).suffix.lower() if fname else ""

            # Image extension set — used below for the OCR / description path
            _EML_ATT_IMG_EXTS = {
                ".png", ".jpg", ".jpeg", ".gif",
                ".bmp", ".tiff", ".tif", ".webp",
            }
            # CSV/TSV extensions handled via PostgreSQL import (separate from _DISPATCH)
            _EML_CSV_EXTS = {".csv", ".tsv"}
            # Min word count from OCR to classify an image as text-heavy.
            # 20 words ≈ 30 tokens — mirrors _IMG_MIN_OCR_TOKENS in chunker.
            _EML_OCR_MIN_WORDS = 20

            if payload and ext in _EML_CSV_EXTS:
                # ── CSV/TSV attachment → PostgreSQL import ────────────────────────
                # Mirror the standalone CSV upload path in api.py:
                #   1. Write payload to a temp file (import_csv needs a real path)
                #   2. Derive a stable table name from the attachment filename
                #      so the same CSV in future emails reuses / updates the table
                #   3. Call import_csv() — hash-tracked, skips if unchanged
                #   4. Emit a schema text block so the CSV is discoverable in
                #      vector search without needing to query the DB
                # Derive a clean PostgreSQL table name from the attachment filename.
                # Replicate csv_executer._slugify() inline to avoid a circular import.
                def _eml_slugify(name: str) -> str:
                    stem = re.sub(r"[^a-z0-9]+", "_", Path(name).stem.lower()).strip("_")
                    return ("t_" + stem) if stem and stem[0].isdigit() else (stem or "csv_att")

                _tbl_name  = _eml_slugify(fname)
                _tmp_path  = None
                _page_offset = 1000 * (len(attachments) + 1)

                try:
                    # Write payload to a named temp file — import_csv() needs a path
                    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as _tmp:
                        _tmp.write(payload)
                        _tmp_path = _tmp.name

                    try:
                        from csv_query_engine.csv_executer import (
                            import_csv    as _csv_import,
                            get_table_info as _csv_table_info,
                        )
                        _csv_available = True
                    except Exception as _csv_imp_err:
                        logger.warning(
                            f"  EML — csv_executer unavailable ({_csv_imp_err}); "
                            f"falling back to plain text extraction for {fname}"
                        )
                        _csv_available = False

                    if _csv_available:
                        logger.info(f"  EML — importing CSV attachment to PostgreSQL: {fname}")
                        _table_name = _csv_import(_tmp_path, table_name=_tbl_name)
                        _info       = _csv_table_info(_table_name)
                        _col_names  = ", ".join(_info["columns"])
                        _nrows      = _info["shape"][0]
                        _ncols      = _info["shape"][1]

                        # Schema text block — makes the CSV searchable in vector index:
                        # "what tables did Majd attach?" → this chunk surfaces it.
                        _schema = (
                            f"[attachment:{fname}]\n"
                            f"CSV data imported into PostgreSQL table '{_table_name}'.\n"
                            f"Rows: {_nrows}  |  Columns ({_ncols}): {_col_names}\n"
                            f"Query this table via the /query endpoint."
                        )
                        blocks.append((_page_offset, _schema))
                        attachments.append({
                            "name":       fname,
                            "extracted":  True,
                            "method":     "postgresql",
                            "table_name": _table_name,
                            "rows":       _nrows,
                            "cols":       _ncols,
                            "columns":    _info["columns"],
                        })
                        logger.info(
                            f"  EML — {fname}: imported → '{_table_name}' "
                            f"({_nrows} rows, {_ncols} cols)"
                        )

                    else:
                        # PostgreSQL not available — fall back to _parse_csv for
                        # text/table extraction so at least the data lands in chunks
                        att_result = _parse_csv(Path(_tmp_path))
                        if len(att_result) == 4:
                            att_blocks, att_tables, att_images, _ = att_result
                        else:
                            att_blocks, att_tables, att_images = att_result

                        for pg, text in att_blocks:
                            blocks.append((_page_offset + pg,
                                           f"[attachment:{fname}]\n{text}"))
                        for tbl in att_tables:
                            tables.append(ExtractedTable(
                                page_number=tbl.page_number + _page_offset,
                                table_index=tbl.table_index,
                                markdown=f"<!-- attachment:{fname} -->\n{tbl.markdown}",
                                raw_rows=tbl.raw_rows,
                            ))
                        attachments.append({
                            "name": fname, "extracted": True,
                            "method": "plain_csv_fallback",
                        })

                except Exception as _csv_exc:
                    logger.warning(f"  EML — CSV attachment failed for {fname}: {_csv_exc}")
                    attachments.append({"name": fname, "extracted": False, "error": str(_csv_exc)})
                finally:
                    if _tmp_path:
                        try: os.unlink(_tmp_path)
                        except Exception: pass

            elif payload and ext in _DISPATCH:
                tmp_path = None
                try:
                    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
                        tmp.write(payload)
                        tmp_path = tmp.name

                    logger.info(f"  EML — extraction pièce jointe : {fname}")
                    att_result = _DISPATCH[ext](Path(tmp_path))

                    # Toutes les parsers retournent un 4-tuple depuis la refonte
                    if len(att_result) == 4:
                        att_blocks, att_tables, att_images, _att_meta = att_result
                    else:
                        att_blocks, att_tables, att_images = att_result

                    # Décalage de page pour éviter la collision avec la page 0 (corps EML)
                    page_offset = 1000 * (len(attachments) + 1)

                    for pg, text in att_blocks:
                        blocks.append((pg + page_offset,
                                       f"[attachment:{fname}]\n{text}"))

                    for tbl in att_tables:
                        tables.append(ExtractedTable(
                            page_number=tbl.page_number + page_offset,
                            table_index=tbl.table_index,
                            markdown=f"<!-- attachment:{fname} -->\n{tbl.markdown}",
                            raw_rows=tbl.raw_rows,
                        ))

                    for img in att_images:
                        new_img = ExtractedImage(
                            page_number=img.page_number + page_offset,
                            image_index=img_idx,
                            base64_data=img.base64_data,
                            mime_type=img.mime_type,
                            caption=(f"[From attachment: {fname}] {img.caption or ''}").strip(),
                        )
                        images.append(new_img)
                        img_idx += 1

                    attachments.append({
                        "name": fname,
                        "extracted": True,
                        "text_blocks": len(att_blocks),
                        "tables": len(att_tables),
                        "images": len(att_images),
                    })
                    logger.info(f"  EML — {fname} : {len(att_blocks)} blocs, "
                                f"{len(att_tables)} tables, {len(att_images)} images")

                except Exception as exc:
                    logger.warning(f"  EML — échec extraction {fname}: {exc}")
                    attachments.append({"name": fname, "extracted": False, "error": str(exc)})
                finally:
                    if tmp_path:
                        try:
                            os.unlink(tmp_path)
                        except Exception:
                            pass
            elif payload and ext in _EML_ATT_IMG_EXTS:
                # ── Image attachment — OCR or Groq description ──────────────────
                # Strategy:
                #   1. Run OCR (Groq / Llama 4 Scout) on the image bytes.
                #   2. If OCR returns ≥ _EML_OCR_MIN_WORDS words → text-heavy image
                #      → add as an attachment text block so the TXT semantic
                #        pipeline can chunk and embed it.
                #   3. Otherwise → regular image (photo, logo, diagram)
                #      → run Groq vision description → store as ExtractedImage
                #        so the chunker emits it as an image chunk with caption.

                # ── Determine MIME type from extension ───────────────────────────
                _mime_ext = ext.lstrip(".")
                if _mime_ext == "jpg":
                    _mime_ext = "jpeg"
                _mime = f"image/{_mime_ext}"

                # ── Step 1: OCR attempt ──────────────────────────────────────────
                _ocr_text = ""
                try:
                    from file_processor.gemma4 import gemma4_available, ocr_image_bytes_gemma
                    if gemma4_available():
                        logger.info(f"  EML — OCR scanning image attachment: {fname}")
                        _ocr_text = ocr_image_bytes_gemma(payload)
                except Exception as _ocr_err:
                    logger.warning(f"  EML — OCR failed for {fname}: {_ocr_err}")

                _ocr_words = len(_ocr_text.split()) if _ocr_text else 0
                page_offset = 1000 * (len(attachments) + 1)

                if _ocr_words >= _EML_OCR_MIN_WORDS:
                    # ── Text-heavy image → text block ────────────────────────────
                    # Add OCR output as an attachment text block.  The chunker's
                    # EML attachment path (page >= 1000) picks it up automatically,
                    # running it through the TXT semantic pipeline with full
                    # email metadata propagation.
                    blocks.append((page_offset, f"[attachment:{fname}]\n{_ocr_text}"))
                    attachments.append({
                        "name":      fname,
                        "extracted": True,
                        "method":    "ocr",
                        "ocr_words": _ocr_words,
                    })
                    logger.info(
                        f"  EML — {fname}: text-heavy image, "
                        f"OCR extracted {_ocr_words} words → text block"
                    )

                else:
                    # ── Regular image → Groq visual description ──────────────────
                    # Generate a caption with Groq's vision model (same model
                    # used for inline image captioning).  The image + caption
                    # land in `images` and are emitted as image chunks by the
                    # chunker with section=email_attachment and email metadata.
                    _b64_data = base64.b64encode(payload).decode()
                    _caption  = ""
                    try:
                        from file_processor.groq_client import caption_image_groq
                        logger.info(
                            f"  EML — generating Groq description for "
                            f"image attachment: {fname}"
                        )
                        _caption = caption_image_groq(_b64_data, mime=_mime)
                    except Exception as _cap_err:
                        logger.warning(
                            f"  EML — Groq description failed for {fname}: {_cap_err}"
                        )

                    _full_caption = (
                        f"[From attachment: {fname}] {_caption}".strip()
                        if _caption else f"[From attachment: {fname}]"
                    )
                    img_obj = ExtractedImage(
                        page_number=page_offset,
                        image_index=img_idx,
                        base64_data=_b64_data,
                        mime_type=_mime,
                        caption=_full_caption,
                    )
                    images.append(img_obj)
                    img_idx += 1
                    attachments.append({
                        "name":      fname,
                        "extracted": True,
                        "method":    "description",
                        "captioned": bool(_caption),
                    })
                    logger.info(
                        f"  EML — {fname}: regular image, "
                        f"Groq description → image chunk (caption: {bool(_caption)})"
                    )

            else:
                # Format non supporté ou payload vide
                reason = "unsupported format" if ext not in _DISPATCH else "no payload"
                attachments.append({"name": fname, "extracted": False, "reason": reason})

    # Attachments in metadata
    if attachments:
        doc_metadata["attachments"] = attachments

    # Corps du message → text blocks uniquement (sans métadonnées)
    # html_blocks fallback: normalise CRLF too before splitting
    if not plain_blocks and html_blocks:
        html_blocks = [b.replace("\r\n", "\n").replace("\r", "\n") for b in html_blocks]
    # Prefer HTML body when it contained tables: the HTML version has the
    # <table> elements replaced with whitespace (clean prose), whereas the
    # text/plain part from Outlook dumps each cell on its own line, producing
    # garbled one-word paragraphs that embed very poorly.
    if plain_blocks and html_blocks and html_had_tables:
        logger.debug("  EML — using HTML body (tables found); discarding text/plain duplicate")
        body_parts = html_blocks
    else:
        body_parts = plain_blocks if plain_blocks else html_blocks
    for text in body_parts:
        # Split on one or more blank lines (handles single \n\n and \n\n\n etc.)
        for para in re.split(r'\n{2,}', text):
            para = para.strip()
            # Unwrap soft line breaks — email clients wrap at ~80 chars using
            # a single \n.  Collapsing these gives the chunker clean sentences
            # instead of mid-sentence line breaks in the chunk text.
            para = re.sub(r'(?<!\n)\n(?!\n)', ' ', para)
            if para:
                blocks.append((0, para))

    logger.info(f"  EML — {len(blocks)} blocs, {len(tables)} tables, "
                f"{len(images)} images, {len(attachments)} pièce(s) jointe(s)")
    return blocks, tables, images, doc_metadata


def _parse_csv(path: Path):
    """
    Extraction CSV/TSV avec parsing tabulaire :
    - Détection automatique du délimiteur (virgule, point-virgule, tabulation, pipe)
    - En-tête détectée automatiquement
    - Résumé texte : nom de fichier, nb lignes, colonnes listées
    - Batches de 50 lignes max (header répété dans chaque batch) → table chunks
    - Même logique que XLSX pour cohérence RAG
    """
    import csv as _csv

    BATCH_SIZE = 50

    with open(path, "rb") as f:
        enc = _encoding(f)

    with open(path, newline="", encoding=enc, errors="replace") as f:
        sample = f.read(8192)

    # Détection du délimiteur
    try:
        dialect = _csv.Sniffer().sniff(sample, delimiters=",;\t|")
        delimiter = dialect.delimiter
    except _csv.Error:
        # Fallback : tabulation si .tsv, sinon virgule
        delimiter = "\t" if path.suffix.lower() == ".tsv" else ","

    rows: list[list[str]] = []
    try:
        with open(path, newline="", encoding=enc, errors="replace") as f:
            reader = _csv.reader(f, delimiter=delimiter)
            for row in reader:
                cleaned = [c.strip() for c in row]
                # Ignorer les lignes entièrement vides
                if any(cleaned):
                    rows.append(cleaned)
    except Exception as e:
        logger.warning(f"  CSV parsing échoué: {e} — fallback texte brut")
        return _parse_text(path)

    if not rows:
        logger.warning(f"  CSV vide ou illisible: {path.name}")
        return [], [], []

    # Détection automatique de l'en-tête
    try:
        has_header = _csv.Sniffer().has_header(sample)
    except _csv.Error:
        has_header = True   # on suppose qu'il y a un header

    if has_header:
        header    = rows[0]
        data_rows = rows[1:]
    else:
        header    = [f"col_{i+1}" for i in range(len(rows[0]))]
        data_rows = rows

    blocks: list[tuple[int, str]] = []
    tables: list[ExtractedTable]  = []

    # Résumé texte (embeddable context pour le fichier entier)
    cols_preview = ", ".join(header[:8]) + ("…" if len(header) > 8 else "")
    summary = (
        f"Fichier CSV '{path.name}': {len(data_rows)} lignes, "
        f"{len(header)} colonnes. "
        f"Colonnes : {cols_preview}. "
        f"Délimiteur : '{delimiter}'."
    )
    blocks.append((0, summary))

    n_batches = max(1, (len(data_rows) + BATCH_SIZE - 1) // BATCH_SIZE)

    if n_batches == 1:
        sheet_prefix = f"CSV: {path.name}\n"
        tables.append(ExtractedTable(
            page_number=0, table_index=0,
            markdown=sheet_prefix + _rows_to_md([header] + data_rows),
            raw_rows=[header] + data_rows,
        ))
    else:
        logger.info(f"  CSV '{path.name}': {len(data_rows)} lignes → {n_batches} batches de {BATCH_SIZE}")
        for b in range(n_batches):
            batch_data = data_rows[b * BATCH_SIZE:(b + 1) * BATCH_SIZE]
            batch_rows = [header] + batch_data
            start_row  = b * BATCH_SIZE + 2
            end_row    = start_row + len(batch_data) - 1
            sheet_prefix = f"CSV: {path.name}, rows {start_row}–{end_row}\n"
            tables.append(ExtractedTable(
                page_number=0, table_index=b,
                markdown=sheet_prefix + _rows_to_md(batch_rows),
                raw_rows=batch_rows,
            ))

    doc_metadata: dict = {
        "file_type":    path.suffix.lower().lstrip("."),
        "delimiter":    delimiter,
        "row_count":    len(data_rows),
        "column_count": len(header),
        "column_names": header,
        "encoding":     enc,
        "has_header":   has_header,
    }

    logger.info(f"  CSV — {len(data_rows)} lignes, {len(header)} colonnes, {len(tables)} chunk(s) table")
    return blocks, tables, [], doc_metadata


def _parse_text(path: Path):
    with open(path, "rb") as f: enc = _encoding(f)
    text = path.read_text(encoding=enc, errors="replace")
    # Regroupe par paragraphes (double saut de ligne) puis par lignes si pas de paragraphes
    raw_blocks = [b.strip() for b in text.split("\n\n") if b.strip()]
    if not raw_blocks:
        raw_blocks = [l.strip() for l in text.splitlines() if l.strip()]
    blocks = [(0, b) for b in raw_blocks]

    doc_metadata: dict = {
        "file_type":  path.suffix.lower().lstrip(".") or "txt",
        "encoding":   enc,
        "line_count": text.count("\n") + 1,
        "word_count": len(text.split()),
        "char_count": len(text),
    }
    return blocks, [], [], doc_metadata


# ── Images scannées (.png .jpg .jpeg .webp) ───────────────────────────────────

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}

def _parse_image(path: Path):
    """
    Traite une image scannée :
    1. Charge et stocke comme ExtractedImage (base64)
    2. Applique l'OCR pour extraire le texte visible
    """
    raw = path.read_bytes()
    img = _make_image(raw, page=0, idx=0)

    if img is None:
        logger.warning(f"  Image trop petite ou invalide — ignorée.")
        return [], [], [], {"file_type": path.suffix.lower().lstrip(".")}

    # ── Métadonnées image ─────────────────────────────────────────────────────
    doc_metadata: dict = {"file_type": path.suffix.lower().lstrip(".")}
    try:
        pil = Image.open(io.BytesIO(raw))
        doc_metadata["width"]       = pil.width
        doc_metadata["height"]      = pil.height
        doc_metadata["format"]      = pil.format or path.suffix.upper().lstrip(".")
        doc_metadata["color_mode"]  = pil.mode
        doc_metadata["file_size_kb"] = round(len(raw) / 1024, 1)
        exif = pil._getexif() if hasattr(pil, "_getexif") else None
        if exif:
            from PIL.ExifTags import TAGS
            for tag_id, val in exif.items():
                tag = TAGS.get(tag_id, tag_id)
                if tag in ("DateTime", "Make", "Model", "Software", "Artist", "Copyright"):
                    doc_metadata[f"exif_{tag}"] = str(val)
    except Exception as e:
        logger.debug(f"  image metadata: {e}")

    text_blocks = []
    from file_processor.gemma4 import gemma4_available, ocr_image_bytes_gemma, OCRMode as _GemmaOCRMode

    # ── Handwriting hint from filename ────────────────────────────────────────
    # Gemma 4 handles all scripts (Arabic, Latin, CJK …) and both handwriting
    # and print automatically — no language flags needed.
    # We use the handwriting hint to select the appropriate OCR mode so Gemma 4
    # applies its handwriting-tuned prompt for those images.
    stem_lower = path.stem.lower()
    _is_handwriting = any(k in stem_lower for k in
                          ("handwrit", "handwritten", "manuscrit", "handwriting"))

    # ── Gemma 4 OCR ───────────────────────────────────────────────────────────
    if gemma4_available():
        logger.info(f"  Gemma 4 OCR on '{path.name}' (handwriting={_is_handwriting}) ...")
        try:
            _ocr_mode = _GemmaOCRMode.HANDWRITING if _is_handwriting else _GemmaOCRMode.EXTRACT
            ocr_text = ocr_image_bytes_gemma(
                raw,
                mode=_ocr_mode,
            )
            if ocr_text.strip():
                # Gemma 4 returns structured markdown; preserve paragraph structure.
                # Split on double-newline so table rows and heading lines stay together.
                paras = [p.strip() for p in ocr_text.split("\n\n") if p.strip()]
                text_blocks = [(0, p) for p in paras]
                logger.info(f"  Gemma 4 OCR: {len(text_blocks)} block(s) extracted")
            else:
                logger.info("  Gemma 4 OCR: no text detected")
        except Exception as e:
            logger.warning(f"  Gemma 4 OCR failed: {e}")

    if not text_blocks:
        logger.info("  No OCR engine available — image stored without text")

    return text_blocks, [], [img], doc_metadata


# ── Dispatch table ─────────────────────────────────────────────────────────────

_DISPATCH: dict[str, Callable] = {
    ".docx": _parse_docx,
    ".pdf":  _parse_pdf,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
    ".eml":  _parse_eml,
}
for _e in PLAIN_TEXT_EXTS:
    _DISPATCH[_e] = _parse_text
for _e in CSV_EXTS:
    _DISPATCH[_e] = _parse_csv
for _e in IMAGE_EXTS:
    _DISPATCH[_e] = _parse_image

SUPPORTED = set(_DISPATCH)


def extract(
    file_path: str | Path,
    pdf_pass: Optional[str] = None,
) -> tuple[list, list[ExtractedTable], list[ExtractedImage], dict]:
    """
    Point d'entrée unique.
    Retourne (text_blocks, tables, images, doc_metadata).
    text_blocks  : liste de (page_number, text)
    tables       : ExtractedTable avec .markdown et .raw_rows
    images       : ExtractedImage avec .base64_data
    doc_metadata : dict de métadonnées document (ex: en-têtes email)
    """
    path = Path(file_path)
    ext  = path.suffix.lower()

    if ext not in _DISPATCH:
        raise ValueError(f"Extension non supportée: '{ext}'. Supportées: {sorted(SUPPORTED)}")

    if ext in PASSWORD_EXTS:
        with open(path, "rb") as f:
            checker = _LOCK_CHECK.get(ext)
            if checker and checker(f):
                logger.warning(f"'{path.name}' est protégé par mot de passe — ignoré.")
                return [], [], [], {"file_type": ext.lstrip("."), "encrypted": True}

    logger.info(f"Extraction de '{path.name}' ...")
    fn = _DISPATCH[ext]
    if ext == ".pdf":
        result = fn(path, pdf_pass=pdf_pass)
    else:
        result = fn(path)

    # All parsers now return 4 elements; fallback to empty dict if somehow 3
    if len(result) == 4:
        return result[0], result[1], result[2], result[3]
    return result[0], result[1], result[2], {"file_type": ext.lstrip(".")}